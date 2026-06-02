import PyPDF2
from pptx import Presentation
from pptx.util import Inches

def pdf_to_text_with_ocr(pdf_path):
    text = ''
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + '\n'
    return text

def create_pptx(text, pptx_path):
    presentation = Presentation()

    # Create a title slide
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "OCR Text Presentation"
    subtitle.text = "Converted text from PDF with OCR"

    # Create content slides
    bullet_slide_layout = presentation.slide_layouts[1]

    for paragraph in text.split('\n'):
        content_slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = content_slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Extracted Text"
        tf = body_shape.text_frame
        tf.text = paragraph

    # Save the presentation
    presentation.save(pptx_path)

if __name__ == "__main__":
    input_pdf_path = 'input.pdf'
    output_pptx_path = 'outputx.pptx'

    extracted_text = pdf_to_text_with_ocr(input_pdf_path)
    create_pptx(extracted_text, output_pptx_path)

    print(f"Conversion completed. PowerPoint presentation saved at: {output_pptx_path}")
